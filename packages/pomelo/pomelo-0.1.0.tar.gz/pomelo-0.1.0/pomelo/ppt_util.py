from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm

class PPTUtil(object):
    def __init__(self,ppt_temp):
        self.prs = Presentation(ppt_temp)

    '''
        新增报表幻灯片
        title ：幻灯片标题
        x_lab: 报表x轴
        x_val: 报表y轴
        slide_idx：依赖哪个幻灯片模版
        title_dix: 标题在模版中的位置
    '''
    def new_chart_slide(self,title,x_lab,x_val,slide_idx = 3,title_idx = 0):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[slide_idx])
        placeholder = slide.placeholders[title_idx]
        placeholder.text = title
        chart_data = CategoryChartData()
        chart_data.categories = x_lab
        chart_data.add_series('',x_val)

        x, y, cx, cy = Inches(0.5), Inches(2), Inches(12), Inches(3.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(6)
        category_axis.has_major_gridlines = False
        category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE

        value_axis = chart.value_axis
        tick_labels = value_axis.tick_labels
        tick_labels.font.size = Pt(8) 
        value_axis.has_major_gridlines = False
        value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE

        bar_plot = chart.plots[0]
        bar_plot.gap_width = 20

        for i in range(len(x_lab)):
            point = chart.series[0].points[i]  # a (data) point in a bar chart is a bar            
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(102,171,163)
        tot = 0
        for v in x_val:
            tot = tot + v

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(6), Inches(6), Inches(1))
        text_frame = shape.text_frame
        text_frame.clear()  # not necessary for newly-created shape
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(tot)
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(70)
        font.bold = True
    
    def new_tab_slide(self,title,tab_head,tab_data,slide_idx = 3,title_idx = 0,col_width = None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[slide_idx])
        placeholder = slide.placeholders[title_idx]
        placeholder.text = title
        top    = Cm(3.5)
        left   = Cm(0.5)
        width  = Cm(28)
        height = Cm(10)

        table = slide.shapes.add_table(len(tab_data) + 1, len(tab_head), left, top, width, height).table
        for i in range(len(tab_head)):
            if (col_width == None or col_width[i] == None):
                table.columns[i].width = Cm(28/len(tab_head))
            else:
                table.columns[i].width = Cm(col_width[i])
        if (col_width == None):
            table.columns[len(tab_head) - 1].width = Cm(10)

        for i,v in enumerate(tab_head):
            cell = table.cell(0,i)
            cell.text = v
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)

        for i,v in enumerate(tab_data):
            for j,d in enumerate(v):
                cell = table.rows[i + 1].cells[j]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                if (j == 7 and float(d) > 0):
                    paragraph.font.color.rgb = RGBColor(255, 0, 0)
                elif (j == 7 and float(d) < 0):
                    paragraph.font.color.rgb = RGBColor(0, 255, 0)
                else:
                    paragraph.font.color.rgb = RGBColor(12, 34, 56)
                cell.text = str(d)        


    def save(self,out_path,end_slide = None):
        if (end_slide != None):
            self.prs.slides.add_slide(self.prs.slide_layouts[end_slide])
        self.prs.save(out_path)